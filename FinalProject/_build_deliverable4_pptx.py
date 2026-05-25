"""Build Deliverable 4 presentation deck (8-minute video)."""
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Deliverable4_Presentation-emmanart.pptx"
PROJECT_DIR = Path(__file__).resolve().parent
ERD_PNG = PROJECT_DIR / "soccer_proj_erd_slide.png"
# If you add your DBeaver export here, the deck will prefer it over the generated diagram.
USER_ERD_CANDIDATES = [
    PROJECT_DIR / "s26a_db0061 - s26a_db0061 - soccer_proj.png",
    PROJECT_DIR / "soccer_proj_erd.png",
]
DEMO_Q16_RESULTS_PNG = PROJECT_DIR / "demo_q16_results.png"
DEMO_Q18_RESULTS_PNG = PROJECT_DIR / "demo_q18_results.png"

GREEN = RGBColor(0x1B, 0x5E, 0x20)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    tf.text = text


def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.font.name = "Calibri"


def add_bullets(slide, items: list[str], top=1.35, left=0.6, width=8.8, height=5.5, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


DEMO_Q16_SQL = """SELECT c.name AS country,
  AVG((m.home_team_goal + m.away_team_goal)::numeric) AS avg_goals_per_match
FROM soccer_proj."Match_tbl" m
JOIN soccer_proj."League" l ON l.id = m.league_id
JOIN soccer_proj."Country" c ON c.id = l.country_id
GROUP BY c.name
ORDER BY avg_goals_per_match DESC;"""

DEMO_Q18_SQL = """SELECT t.team_long_name,
  SUM(CASE WHEN m.home_team_api_id = t.team_api_id AND m.away_team_goal = 0 THEN 1 ELSE 0 END
     + CASE WHEN m.away_team_api_id = t.team_api_id AND m.home_team_goal = 0 THEN 1 ELSE 0 END) AS clean_sheets
FROM soccer_proj."Team" t
JOIN soccer_proj."Match_tbl" m
  ON m.home_team_api_id = t.team_api_id OR m.away_team_api_id = t.team_api_id
GROUP BY t.team_long_name
ORDER BY clean_sheets DESC
LIMIT 30;"""


def add_result_screenshot(slide, image_path: Path, top=1.1, left=0.35, max_width=9.3, max_height=6.15):
    """Scale DBeaver screenshot to fit below slide title."""
    if not image_path.is_file():
        add_bullets(
            slide,
            [f"[Missing screenshot: {image_path.name}]"],
            top=top + 0.5,
            size=18,
        )
        return
    pic = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(max_width))
    if pic.height > Inches(max_height):
        scale = Inches(max_height) / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(Inches(left))
        pic.top = int(Inches(top))


def add_demo_results_slide(prs, title: str, subtitle: str, image_path: Path, notes: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, title, subtitle)
    add_result_screenshot(s, image_path)
    set_notes(s, notes)


def add_sql_block(slide, sql: str, top=3.55, left=0.55, width=8.9, height=3.5, size=9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SQL (copy into DBeaver on soccer_proj):"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    p2.text = sql.strip()
    p2.font.size = Pt(size)
    p2.font.name = "Consolas"
    p2.font.color.rgb = DARK
    p2.space_before = Pt(4)


def _draw_entity(
    ax,
    x,
    y,
    w,
    h,
    title,
    attrs,
    *,
    header="#1B5E20",
    body="#F1F8F4",
    accent=None,
):
    """Draw a table box: green header + left-aligned attributes."""
    if accent:
        body = accent
    outer = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.01,rounding_size=0.06",
        linewidth=2.2,
        edgecolor=header,
        facecolor="white",
    )
    ax.add_patch(outer)
    header_h = min(0.42, h * 0.32)
    hdr = FancyBboxPatch(
        (x + 0.04, y + h - header_h - 0.04),
        w - 0.08,
        header_h,
        boxstyle="round,pad=0.01,rounding_size=0.04",
        linewidth=0,
        facecolor=header,
    )
    ax.add_patch(hdr)
    ax.text(
        x + w / 2,
        y + h - header_h / 2 - 0.04,
        title,
        ha="center",
        va="center",
        fontsize=12,
        fontweight="bold",
        color="white",
    )
    body_patch = FancyBboxPatch(
        (x + 0.06, y + 0.06),
        w - 0.12,
        h - header_h - 0.12,
        boxstyle="round,pad=0.01,rounding_size=0.03",
        linewidth=0,
        facecolor=body,
    )
    ax.add_patch(body_patch)
    line_h = 0.22
    top_y = y + h - header_h - 0.22
    for i, line in enumerate(attrs):
        ax.text(
            x + 0.18,
            top_y - i * line_h,
            f"• {line}",
            ha="left",
            va="top",
            fontsize=9.5,
            color="#212121",
        )


def _link(ax, start, end, *, dashed=False):
    color = "#616161" if dashed else "#1B5E20"
    style = (0, (5, 4)) if dashed else "-"
    arr = FancyArrowPatch(
        start,
        end,
        arrowstyle="-|>",
        mutation_scale=16,
        linewidth=2.4,
        linestyle=style,
        color=color,
        shrinkA=2,
        shrinkB=2,
        connectionstyle="arc3,rad=0.0",
    )
    ax.add_patch(arr)


def generate_erd_png(path: Path) -> Path:
    """Clean ER diagram for the presentation slide."""
    fig, ax = plt.subplots(figsize=(14, 8.5), dpi=160)
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 9)
    ax.axis("off")
    fig.patch.set_facecolor("white")

    # Center spine: Country → League → Match_tbl
    _draw_entity(ax, 5.35, 7.0, 3.3, 1.05, "Country", ["id (primary key)", "name"])
    _draw_entity(
        ax,
        4.95,
        5.35,
        4.1,
        1.15,
        "League",
        ["id (primary key)", "country_id (foreign key)", "name"],
    )
    _draw_entity(
        ax,
        3.85,
        2.85,
        6.3,
        1.55,
        "Match_tbl",
        [
            "id (primary key)",
            "country_id, league_id (foreign keys)",
            "home / away team (foreign keys)",
            "goals, date, odds, lineups…",
        ],
        accent="#FFFDE7",
    )

    # Left column: Team stack
    _draw_entity(
        ax,
        0.55,
        3.15,
        3.0,
        1.15,
        "Team",
        ["id (primary key)", "team_api_id (unique key)", "team_long_name"],
    )
    _draw_entity(
        ax,
        0.55,
        1.35,
        3.0,
        1.05,
        "Team_Attributes",
        ["id (primary key)", "team_api_id", "date, tactics…"],
    )

    # Right column: Player stack
    _draw_entity(
        ax,
        10.45,
        3.15,
        3.0,
        1.15,
        "Player",
        ["id (primary key)", "player_api_id", "player_name"],
    )
    _draw_entity(
        ax,
        10.15,
        1.25,
        3.35,
        1.1,
        "Player_Attributes",
        ["id (primary key)", "player_api_id", "date, overall_rating…"],
        accent="#E3F2FD",
    )

    # Arrows — no text on lines (legend below)
    _link(ax, (6.95, 7.0), (6.95, 6.5))  # Country → League
    _link(ax, (6.95, 5.35), (6.95, 4.4))  # League → Match
    _link(ax, (3.55, 3.72), (3.85, 3.72))  # Team → Match
    _link(ax, (2.05, 3.15), (2.05, 2.4))  # Team → Team_Attributes
    _link(ax, (11.95, 3.15), (11.82, 2.35), dashed=True)  # Player → Player_Attributes
    _link(ax, (5.35, 7.5), (4.6, 4.4), dashed=True)  # Country → Match (logical)

    # Title + legend (no crowded labels on arrows)
    ax.text(
        0.5,
        8.55,
        "soccer_proj — entity-relationship diagram (7 base tables)",
        fontsize=16,
        fontweight="bold",
        color="#1B5E20",
    )
    legend_lines = [
        "Solid arrow = enforced foreign key in PostgreSQL",
        "Dashed arrow = logical link in source data (may not be enforced)",
        "Views (not drawn): league_match_summary_view, match_context_view",
    ]
    for i, line in enumerate(legend_lines):
        ax.text(0.5, 0.72 - i * 0.28, line, fontsize=10, color="#424242")

    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white", pad_inches=0.2)
    plt.close(fig)
    return path


def resolve_erd_image() -> Path:
    for candidate in USER_ERD_CANDIDATES:
        if candidate.is_file():
            return candidate
    return generate_erd_png(ERD_PNG)  # regenerated each build when no user PNG present


def add_erd_slide(prs: Presentation, erd_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "ER diagram — soccer_proj", "Initial and final use the same 7 tables")

    slide.shapes.add_picture(
        str(erd_path),
        Inches(0.25),
        Inches(1.05),
        width=Inches(9.5),
    )

    cap = slide.shapes.add_textbox(Inches(0.4), Inches(6.45), Inches(9.2), Inches(0.85))
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "Data model for business reporting: Country/League/Team/Player = who & where; "
        "Match_tbl = transactions (each game); attributes = snapshots over time. Views = ready-made reports."
    )
    p.font.size = Pt(13)
    p.font.name = "Calibri"
    p.font.color.rgb = DARK

    set_notes(
        slide,
        "~1.5 min. Walk the diagram: Country → League → Match_tbl; Team for home/away; "
        "attribute history tables. Say final ERD matches initial; views sit on top.",
    )


def add_code_block(slide, code: str, top=3.0):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.9), Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(8.8), Inches(1.4))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "European Soccer Database"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = "Final Project · Deliverable 4 · soccer_proj"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"

    sub = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(8.8), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Emmanuel Arthur · Database Class · Spring 2026"
    sp.font.size = Pt(22)
    sp.font.color.rgb = DARK
    sp.font.name = "Calibri"

    set_notes(
        slide,
        "~30 sec. Introduce yourself and the project. Mention Kaggle European Soccer data, "
        "Postgres schema soccer_proj, and that the video covers design, reflection, and a short demo.",
    )


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    title_slide(prs)

    # --- BEGINNING: overview, ERD, reflection (per instructions 35-39, 45-53) ---

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Dataset overview")
    add_bullets(
        s,
        [
            "Source: Kaggle European Soccer Database (SQLite export → CSV → PostgreSQL)",
            "Scope: 11 countries, 11 leagues, ~26K matches (2008–2016), teams, players, attributes",
            "Business context: same kind of data clubs, sports media, and analytics vendors use for reporting",
            "Schema soccer_proj with 7 base tables + views for common joins",
            "Total loaded rows across all tables: 222,796",
        ],
        size=18,
    )
    set_notes(
        s,
        "~1 min. Explain scope, then one sentence on business: organizations use match + player history "
        "for performance review, content strategy, and market comparison—not just hobby stats.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Why I chose this dataset")
    add_bullets(
        s,
        [
            "I have followed European club soccer since childhood (Man United, then Bayern Munich)",
            "The domain is familiar: leagues, seasons, home/away, points, streaks, betting odds",
            "Business-style questions: league scoring trends, defensive strength, player development, odds vs results",
            "Good mix of dimensions (country, league, team, player) and facts (matches, ratings)",
            "Graduate requirement: 30 English questions that map to real SELECT workloads",
        ],
        size=18,
    )
    set_notes(
        s,
        "~45 sec. Personal tie-in, then say my 30 questions mirror reports a front office or analyst might ask.",
    )

    erd_path = resolve_erd_image()
    add_erd_slide(prs, erd_path)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Hardest part")
    add_bullets(
        s,
        [
            "Loading and typing: SQLite → CSV → Postgres with consistent column names and NULL handling",
            "Match_tbl is very wide (lineup API IDs, many bookmaker columns) — easy to get joins wrong",
            "Writing correct multi-step SQL (common table expressions, window functions) for streaks and lineup-based ratings",
            "Deciding which foreign keys to enforce vs document-only when source data has orphan keys",
            "Keeping 30 project questions aligned with what the data actually supports",
        ],
    )
    set_notes(
        s,
        "~1 min. Be specific: mention one query that was hard (e.g. Q10 unbeaten streak) or one load issue.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "What I learned")
    add_bullets(
        s,
        [
            "Design first: ERD and keys matter before writing dozens of queries",
            "JOIN choice changes results (INNER vs LEFT) — always check row counts",
            "Views encapsulate repeated 4–5 table joins; indexes help filter/join columns on Match_tbl",
            "Window functions (partition rows per team/season, compare to prior rows) for streaks and trends",
            "Documentation + reproducible load scripts save time at deliverable deadline",
        ],
    )
    set_notes(
        s,
        "~1 min. Tie learning to class topics: normalization, SQL, physical design (indexes for grad).",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "SQL confidence now")
    add_bullets(
        s,
        [
            "Comfortable with the most common SQL: SELECT, JOINs, GROUP BY, HAVING, subqueries, CASE",
            "Comfortable creating views (for example league_match_summary_view, match_context_view)",
            "Overall: much more confident than at the start of the semester — I can debug by counting rows",
            "DBeaver plus the EXPLAIN command helped me see when a query execution plan looked expensive",
        ],
    )
    set_notes(
        s,
        "~45 sec. Honest self-assessment; mention you can read and adapt SQL you did not write before.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "What I would do next (more time)")
    add_bullets(
        s,
        [
            "Enforce more foreign keys after data cleaning; add CHECK constraints on goals and dates",
            "Materialized views or summary tables for league/season dashboards",
            "More indexes tuned using EXPLAIN on the slowest of the 30 queries",
            "Data quality report: missing lineups, duplicate player identifiers, odds outliers",
            "Optional: BI dashboard (Tableau/Power BI) on views for non-technical stakeholders",
        ],
        size=18,
    )
    set_notes(
        s,
        "~45 sec. Tie next steps to business: cleaner data + dashboards = decisions faster for ops teams.",
    )

    # --- END: demo slides (instructions 41-43) ---

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "DEMO: Tables + row counts", "DBeaver · schema soccer_proj")
    add_bullets(
        s,
        [
            "Demo these core tables (foreign key relationships) in DBeaver — columns + row count:",
            "Country — 11 rows (parent; League.country_id → Country.id)",
            "League — 11 rows (country_id foreign key; Match_tbl.league_id → League.id)",
            "Team — 299 rows (home/away team_api_id on Match_tbl)",
            "Match_tbl — 25,979 rows (hub: country_id, league_id, home/away teams, goals)",
            "All 7 tables combined — 222,796 rows total",
        ],
        size=18,
    )
    set_notes(
        s,
        "~1 min. LIVE DEMO. Expand Country → League → Match_tbl in the tree. Run SELECT COUNT(*) "
        "on each (or show counts from write-up). Open Match_tbl LIMIT 5 to show league_id, "
        "country_id, home_team_goal, away_team_goal. Mention Player_Attributes is largest "
        "(183,978) but skip opening it unless you have time.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(
        s,
        "DEMO: Intermediate query (Q16)",
        "Joins + GROUP BY + AVG (no CASE)",
    )
    add_bullets(
        s,
        [
            "Question: Which countries had the highest average goals per match?",
            "Tables: Match_tbl → League → Country · INNER JOIN, GROUP BY, AVG (no CASE)",
        ],
        top=1.25,
        height=2.1,
        size=16,
    )
    add_sql_block(s, DEMO_Q16_SQL)
    set_notes(
        s,
        "~1 min. LIVE DEMO. Run Q16.\n\n" + DEMO_Q16_SQL + "\n\n"
        "Say: regional report—which country’s leagues produce the most goals per match.",
    )

    add_demo_results_slide(
        prs,
        "DEMO Q16 — Query results",
        "DBeaver · demo1 · Netherlands leads avg goals per match",
        DEMO_Q16_RESULTS_PNG,
        "~30 sec. Show this screenshot or live results. Highlight top countries (Netherlands, Switzerland, Germany).",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "DEMO: Advanced query (Q18)", "Join with OR + conditional aggregates")
    add_bullets(
        s,
        [
            "Question: Which teams most frequently kept clean sheets?",
            "Tables: Team + Match_tbl · JOIN with OR · two CASE inside SUM · GROUP BY",
        ],
        top=1.25,
        height=2.1,
        size=16,
    )
    add_sql_block(s, DEMO_Q18_SQL, size=8)
    set_notes(
        s,
        "~1 min. LIVE DEMO. Run Q18.\n\n" + DEMO_Q18_SQL + "\n\n"
        "Clean sheets = defensive KPI clubs track home and away.",
    )

    add_demo_results_slide(
        prs,
        "DEMO Q18 — Query results",
        "DBeaver · demo2 · Celtic leads clean sheets",
        DEMO_Q18_RESULTS_PNG,
        "~30 sec. Show screenshot or live results. Celtic, Barcelona, Manchester United at top.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Business applications", "Why this project matters outside class")
    add_bullets(
        s,
        [
            "Club / league operations: season KPIs (goals, clean sheets, home vs away) for coaching and scouting",
            "Media & broadcasting: compare markets (e.g. avg goals by country) for scheduling and storytelling",
            "Betting & risk (data in Match_tbl): historical odds vs outcomes; favorite win rates",
            "Talent / HR analog: Player_Attributes over time ≈ employee skill ratings for hiring and development",
            "Same database patterns as retail (product sales by region) or finance (metrics over time)",
        ],
        size=17,
    )
    set_notes(
        s,
        "~45 sec. Closing slide before Q&A. Pick two examples. Star schema: dimensions + facts; "
        "views = canned reports for business users.",
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Thank you")
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Deliverable 4 write-up + video link in submission"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    set_notes(s, "~15 sec. Close; remind grader where video URL lives in the doc.")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
